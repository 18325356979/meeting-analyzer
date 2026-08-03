#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
会议复盘分析自动化工具1.0 - Streamlit网页版 v2.0
新增功能：历史对比 + 手机适配
"""

import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import io
import os
import json
import tempfile
import platform
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
import warnings
warnings.filterwarnings('ignore')

# ========== 页面配置 ==========
st.set_page_config(
    page_title="会议复盘分析自动化工具1.0",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ========== 中文字体设置 ==========
system = platform.system()
if system == 'Windows':
    plt.rcParams['font.sans-serif'] = ['SimHei', 'Microsoft YaHei']
elif system == 'Darwin':  # macOS
    plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei']
else:  # Linux (Streamlit Cloud 服务器)
    plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'WenQuanYi Zen Hei', 'SimHei']

plt.rcParams['axes.unicode_minus'] = False  # 解决负号显示问题
plt.rcParams['figure.dpi'] = 120

# ========== 历史数据管理 ==========
HISTORY_FILE = os.path.join(os.path.expanduser("~"), ".meeting_analyzer_history.json")

def load_history():
    """加载历史分析记录"""
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return []
    return []

def save_history(record):
    """保存分析记录到历史"""
    history = load_history()
    history.append(record)
    # 只保留最近20条
    history = history[-20:]
    with open(HISTORY_FILE, 'w', encoding='utf-8') as f:
        json.dump(history, f, ensure_ascii=False, indent=2)

def delete_history(index):
    """删除指定历史记录"""
    history = load_history()
    if 0 <= index < len(history):
        history.pop(index)
        with open(HISTORY_FILE, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
        return True
    return False

# ========== CSS样式（优化手机适配） ==========
st.markdown("""
<style>
    .main-title { font-size: 32px; font-weight: bold; color: #1a237e; text-align: center; margin-bottom: 8px; }
    .sub-title { font-size: 14px; color: #666; text-align: center; margin-bottom: 20px; }
    .metric-card { padding: 16px; border-radius: 12px; text-align: center; color: white; margin-bottom: 10px; }
    .metric-value { font-size: 28px; font-weight: bold; }
    .metric-label { font-size: 12px; opacity: 0.9; }
    .stButton>button { background: linear-gradient(135deg, #1a237e 0%, #3949ab 100%); color: white; font-size: 18px; padding: 14px 20px; border-radius: 10px; border: none; width: 100%; font-weight: bold; }
    .stButton>button:hover { background: linear-gradient(135deg, #3949ab 0%, #1a237e 100%); }
    .stDownloadButton>button { background: linear-gradient(135deg, #27ae60 0%, #2ecc71 100%) !important; color: white !important; font-size: 16px !important; padding: 12px 20px !important; border-radius: 10px !important; border: none !important; width: 100% !important; }
    .stDownloadButton>button:hover { background: linear-gradient(135deg, #2ecc71 0%, #27ae60 100%) !important; }
    .info-box { background: #e8f4f8; padding: 12px; border-radius: 8px; border-left: 4px solid #3498db; margin: 10px 0; }
    .warning-box { background: #fdf2e9; padding: 12px; border-radius: 8px; border-left: 4px solid #e67e22; margin: 10px 0; }
    .danger-box { background: #fdedec; padding: 12px; border-radius: 8px; border-left: 4px solid #e74c3c; margin: 10px 0; }
    @media (max-width: 768px) {
        .main-title { font-size: 24px !important; }
        .metric-value { font-size: 22px !important; }
        .stButton>button { font-size: 16px !important; padding: 12px !important; }
    }
</style>
""", unsafe_allow_html=True)


# ========== 核心分析函数 ==========

def parse_duration(duration_str):
    """解析时长字符串为分钟"""
    try:
        duration_str = str(duration_str).strip()
        total_min = 0
        if '小时' in duration_str:
            parts = duration_str.split('小时')
            total_min += int(parts[0]) * 60
            duration_str = parts[1]
        if '分钟' in duration_str:
            parts = duration_str.split('分钟')
            total_min += int(parts[0])
            duration_str = parts[1]
        if '秒' in duration_str:
            sec_str = duration_str.replace('秒', '').strip()
            if sec_str:
                total_min += int(sec_str) / 60
        return round(total_min, 2)
    except:
        return 0

def analyze_meeting(watch_bytes, client_bytes, meeting_name, meeting_date):
    """执行会议分析"""

    # 读取Excel
    watch_df = pd.read_excel(io.BytesIO(watch_bytes))
    client_df = pd.read_excel(io.BytesIO(client_bytes))

    # 清洗客户表
    client_df = client_df.iloc[1:].copy()
    client_df.columns = ['序号', '团队', '专员姓名', '科室', '医生姓名', '医院等级', 
                         '医院', '职称', '梯队', '客户分级', '每月拜访次数', 
                         '季度讲课次数', '维护型', '合作型', '伙伴型', '预估患者总数',
                         '优势患者1', '优势患者1占比', '优势患者2', '优势患者2占比',
                         '当前观念分级', '目标观念升级', '核心竞品']
    client_df = client_df[client_df['医生姓名'].notna() & (client_df['医生姓名'] != '医生姓名')].copy()
    client_df['客户分级'] = client_df['客户分级'].astype(str).str.strip()
    client_df['医生姓名'] = client_df['医生姓名'].astype(str).str.strip()
    client_df['医院'] = client_df['医院'].astype(str).str.strip()

    # 清洗观看记录
    watch_df['用户'] = watch_df['用户'].astype(str).str.strip()
    watch_df['医院'] = watch_df['医院'].astype(str).str.strip()

    def extract_name(user_str):
        user_str = str(user_str).strip()
        if ' ' in user_str:
            return user_str.split()[-1]
        return user_str

    watch_df['用户_纯名'] = watch_df['用户'].apply(extract_name)
    watch_df['用户_纯名'] = watch_df['用户_纯名'].str.replace('。', '', regex=False)

    # 创建客户字典
    client_dict = {}
    for _, row in client_df.iterrows():
        name = str(row['医生姓名']).strip()
        if name not in client_dict:
            client_dict[name] = []
        client_dict[name].append({
            '医院': str(row['医院']).strip(),
            '客户分级': str(row['客户分级']).strip(),
            '团队': row['团队'],
            '专员姓名': row['专员姓名'],
            '科室': row['科室'],
            '职称': row['职称'],
            '医院等级': row['医院等级'],
            '预估患者总数': row['预估患者总数'],
            '当前观念分级': row['当前观念分级'],
            '核心竞品': row['核心竞品'],
        })

    # 匹配
    seen = set()
    results = []
    for _, row in watch_df.iterrows():
        name = row['用户_纯名']
        hospital = row['医院']
        if name in seen:
            continue
        seen.add(name)

        info = None
        grade = '未匹配'
        if name in client_dict:
            entries = client_dict[name]
            if len(entries) == 1:
                info = entries[0]
                grade = info['客户分级']
            else:
                for e in entries:
                    if e['医院'] in hospital or hospital in e['医院']:
                        info = e
                        grade = e['客户分级']
                        break
                if not info:
                    info = entries[0]
                    grade = f"同名({len(entries)})"

        duration_min = parse_duration(row['持续时长'])

        results.append({
            '姓名': name, '医院': hospital, '客户分级': grade,
            '团队': info['团队'] if info else '-',
            '专员': info['专员姓名'] if info else '-',
            '科室': info['科室'] if info else '-',
            '职称': info['职称'] if info else '-',
            '医院等级': info['医院等级'] if info else '-',
            '时长分钟': duration_min,
            '时长文本': row['持续时长'],
            '预估患者总数': info['预估患者总数'] if info else '-',
            '当前观念分级': info['当前观念分级'] if info else '-',
            '核心竞品': info['核心竞品'] if info else '-',
            '是否AB客户': '是' if info else '否',
        })

    result_df = pd.DataFrame(results)

    # 统计
    total = len(result_df)
    a_count = len(result_df[result_df['客户分级'] == 'A'])
    b_count = len(result_df[result_df['客户分级'] == 'B'])
    other = len(result_df[result_df['客户分级'] == '未匹配'])
    total_a_db = len(client_df[client_df['客户分级'] == 'A'])
    total_b_db = len(client_df[client_df['客户分级'] == 'B'])

    stats = {
        'total_attendees': total, 'a_attendees': a_count, 'b_attendees': b_count,
        'other_attendees': other, 'ab_attendees': a_count + b_count,
        'total_a_db': total_a_db, 'total_b_db': total_b_db,
        'coverage_a': round(a_count / total_a_db * 100, 1) if total_a_db > 0 else 0,
        'coverage_b': round(b_count / total_b_db * 100, 1) if total_b_db > 0 else 0,
        'coverage_ab': round((a_count + b_count) / (total_a_db + total_b_db) * 100, 1),
    }

    # 团队数据（用于历史对比）
    ab_data = result_df[result_df['客户分级'].isin(['A', 'B'])]
    team_stats = {}
    if len(ab_data) > 0:
        team_all = client_df.groupby('团队')['客户分级'].value_counts().unstack(fill_value=0)
        team_attend = ab_data.groupby('团队')['客户分级'].value_counts().unstack(fill_value=0) if len(ab_data) > 0 else pd.DataFrame()
        for team in team_all.index:
            total_ab = team_all.loc[team].get('A', 0) + team_all.loc[team].get('B', 0)
            attend_ab = team_attend.loc[team].sum() if team in team_attend.index else 0
            team_stats[team] = {
                'total_ab': int(total_ab),
                'attend_ab': int(attend_ab),
                'coverage': round(attend_ab/total_ab*100, 1) if total_ab > 0 else 0
            }

    # 时长分布
    high_quality = len(result_df[result_df['时长分钟'] > 30])
    low_quality = len(result_df[result_df['时长分钟'] < 5])

    # 保存历史记录
    history_record = {
        'meeting_name': meeting_name,
        'meeting_date': meeting_date,
        'timestamp': datetime.now().isoformat(),
        'stats': {k: v for k, v in stats.items()},
        'team_stats': team_stats,
        'quality': {
            'high': high_quality,
            'medium': len(result_df[(result_df['时长分钟'] >= 5) & (result_df['时长分钟'] <= 30)]),
            'low': low_quality
        }
    }
    save_history(history_record)

    # 生成图表
    fig1 = generate_chart1(result_df, stats, meeting_name)
    fig2 = generate_chart2(result_df, client_df, stats)

    # 生成Excel
    excel_bytes = generate_excel(result_df, client_df, stats)

    # 生成PPT
    ppt_bytes = generate_ppt(result_df, client_df, stats, meeting_name, meeting_date, fig1, fig2)

    return {
        'result_df': result_df,
        'client_df': client_df,
        'stats': stats,
        'fig1': fig1,
        'fig2': fig2,
        'excel_bytes': excel_bytes,
        'ppt_bytes': ppt_bytes,
        'team_stats': team_stats,
        'history_record': history_record,
    }


def generate_chart1(result_df, stats, meeting_name):
    """生成图表1：数据分析看板"""
    fig, axes = plt.subplots(2, 3, figsize=(16, 10))
    fig.suptitle(f'{meeting_name} - 数据分析看板', fontsize=16, fontweight='bold', y=1.02)

    # 饼图
    ax1 = axes[0, 0]
    labels = ['A级客户', 'B级客户', '其他']
    sizes = [stats['a_attendees'], stats['b_attendees'], stats['other_attendees']]
    colors = ['#FF6B6B', '#4ECDC4', '#95A5A6']
    ax1.pie(sizes, explode=(0.05, 0.02, 0), labels=labels, colors=colors, autopct='%1.1f%%',
            shadow=True, startangle=90, textprops={'fontsize': 10})
    ax1.set_title(f'参会构成（{stats["total_attendees"]}人）', fontsize=12, fontweight='bold')

    # 覆盖率对比
    ax2 = axes[0, 1]
    categories = ['A级客户', 'B级客户']
    should = [stats['total_a_db'], stats['total_b_db']]
    actual = [stats['a_attendees'], stats['b_attendees']]
    x = np.arange(len(categories))
    width = 0.35
    ax2.bar(x - width/2, should, width, label='应到', color='#3498DB', alpha=0.8)
    ax2.bar(x + width/2, actual, width, label='实到', color='#E74C3C', alpha=0.8)
    ax2.set_ylabel('人数')
    ax2.set_title('AB级客户覆盖率', fontsize=12, fontweight='bold')
    ax2.set_xticks(x)
    ax2.set_xticklabels(categories)
    ax2.legend(fontsize=9)
    for i, (s, a) in enumerate(zip(should, actual)):
        ax2.text(i - width/2, s + 2, str(s), ha='center', fontsize=9)
        ax2.text(i + width/2, a + 2, str(a), ha='center', fontsize=9)
    ax2.text(0, max(should)*0.5, f'{stats["coverage_a"]}%', ha='center', fontsize=11, color='#E74C3C', fontweight='bold')
    ax2.text(1, max(should)*0.5, f'{stats["coverage_b"]}%', ha='center', fontsize=11, color='#E74C3C', fontweight='bold')

    # 时长分布
    ax3 = axes[0, 2]
    duration_bins = [0, 5, 15, 30, 60, 120]
    duration_labels = ['<5分', '5-15分', '15-30分', '30-60分', '>60分']
    result_df['时长区间'] = pd.cut(result_df['时长分钟'], bins=duration_bins, labels=duration_labels, right=False)
    duration_counts = result_df['时长区间'].value_counts().sort_index()
    colors_bar = ['#E74C3C', '#F39C12', '#F1C40F', '#2ECC71', '#27AE60']
    ax3.bar(duration_counts.index, duration_counts.values, color=colors_bar, edgecolor='white')
    ax3.set_xlabel('观看时长')
    ax3.set_ylabel('人数')
    ax3.set_title('观看时长分布', fontsize=12, fontweight='bold')
    for i, v in enumerate(duration_counts.values):
        ax3.text(i, v + 0.3, str(v), ha='center', fontsize=9, fontweight='bold')

    # 团队分布
    ax4 = axes[1, 0]
    ab_data = result_df[result_df['客户分级'].isin(['A', 'B'])]
    if len(ab_data) > 0:
        team_data = ab_data.groupby(['团队', '客户分级']).size().unstack(fill_value=0)
        team_data = team_data.sort_values('A' if 'A' in team_data.columns else team_data.columns[0], ascending=False)
        team_data.plot(kind='barh', stacked=True, ax=ax4, color=['#FF6B6B', '#4ECDC4'], width=0.7)
        ax4.set_xlabel('人数')
        ax4.set_title('AB级客户按团队分布', fontsize=12, fontweight='bold')
        ax4.legend(['A级', 'B级'], fontsize=9)

    # 医院等级
    ax5 = axes[1, 1]
    if len(ab_data) > 0:
        hosp_level = ab_data['医院等级'].value_counts()
        colors_hl = ['#9B59B6', '#3498DB', '#1ABC9C']
        ax5.pie(hosp_level.values, labels=hosp_level.index, colors=colors_hl, autopct='%1.1f%%',
                shadow=True, startangle=90, textprops={'fontsize': 10})
        ax5.set_title('医院等级分布', fontsize=12, fontweight='bold')

    # 职称
    ax6 = axes[1, 2]
    if len(ab_data) > 0:
        title_data = ab_data['职称'].value_counts()
        ax6.barh(title_data.index, title_data.values, color=['#E67E22', '#D35400', '#C0392B'], height=0.6)
        ax6.set_xlabel('人数')
        ax6.set_title('职称分布', fontsize=12, fontweight='bold')
        for i, v in enumerate(title_data.values):
            ax6.text(v + 0.2, i, str(v), va='center', fontsize=9, fontweight='bold')

    plt.tight_layout()
    return fig

def generate_chart2(result_df, client_df, stats):
    """生成图表2：深度分析"""
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    fig.suptitle('深度分析 - 覆盖率与竞品', fontsize=16, fontweight='bold', y=1.02)

    # 团队覆盖率
    ax1 = axes[0, 0]
    team_all = client_df.groupby('团队')['客户分级'].value_counts().unstack(fill_value=0)
    ab_data = result_df[result_df['客户分级'].isin(['A', 'B'])]
    team_attend = ab_data.groupby('团队')['客户分级'].value_counts().unstack(fill_value=0) if len(ab_data) > 0 else pd.DataFrame()

    team_coverage = []
    for team in team_all.index:
        total_ab = team_all.loc[team].get('A', 0) + team_all.loc[team].get('B', 0)
        attend_ab = team_attend.loc[team].sum() if team in team_attend.index else 0
        team_coverage.append({'团队': team, '覆盖率': attend_ab/total_ab*100 if total_ab > 0 else 0})

    team_cov_df = pd.DataFrame(team_coverage).sort_values('覆盖率', ascending=True)
    colors_cov = ['#E74C3C' if x < 30 else '#F39C12' if x < 50 else '#27AE60' for x in team_cov_df['覆盖率']]
    ax1.barh(team_cov_df['团队'], team_cov_df['覆盖率'], color=colors_cov, height=0.6)
    ax1.set_xlabel('覆盖率 (%)')
    ax1.set_title('各团队AB级客户参会覆盖率', fontsize=12, fontweight='bold')
    ax1.axvline(x=50, color='red', linestyle='--', alpha=0.5)

    # 时长箱线图
    ax2 = axes[0, 1]
    a_dur = result_df[result_df['客户分级'] == 'A']['时长分钟'].values
    b_dur = result_df[result_df['客户分级'] == 'B']['时长分钟'].values
    o_dur = result_df[result_df['客户分级'] == '未匹配']['时长分钟'].values
    bp = ax2.boxplot([a_dur, b_dur, o_dur], tick_labels=['A级', 'B级', '其他'],
                     patch_artist=True, widths=0.5)
    for patch, color in zip(bp['boxes'], ['#FF6B6B', '#4ECDC4', '#95A5A6']):
        patch.set_facecolor(color)
        patch.set_alpha(0.7)
    ax2.set_ylabel('观看时长（分钟）')
    ax2.set_title('观看时长分布（按客户分级）', fontsize=12, fontweight='bold')
    ax2.grid(axis='y', alpha=0.3)

    # 竞品
    ax3 = axes[1, 0]
    if len(ab_data) > 0:
        comp_data = ab_data['核心竞品'].value_counts().head(8)
        colors_comp = plt.cm.Set3(np.linspace(0, 1, len(comp_data)))
        ax3.pie(comp_data.values, labels=comp_data.index, colors=colors_comp, autopct='%1.1f%%',
                shadow=True, startangle=45, textprops={'fontsize': 8})
        ax3.set_title('核心竞品分布', fontsize=12, fontweight='bold')

    # 未参会
    ax4 = axes[1, 1]
    all_ab_names = set(client_df[client_df['客户分级'].isin(['A','B'])]['医生姓名'])
    attend_names = set(ab_data['姓名'].tolist()) if len(ab_data) > 0 else set()
    not_attend = all_ab_names - attend_names
    not_attend_df = client_df[client_df['医生姓名'].isin(not_attend)]
    if len(not_attend_df) > 0:
        not_team = not_attend_df.groupby('团队').size().sort_values(ascending=True)
        ax4.barh(not_team.index, not_team.values, color='#E74C3C', alpha=0.8, height=0.6)
        ax4.set_xlabel('未参会AB客户人数')
        ax4.set_title('各团队未参会AB级客户数', fontsize=12, fontweight='bold')
        for i, v in enumerate(not_team.values):
            ax4.text(v + 0.5, i, str(v), va='center', fontsize=9, fontweight='bold')

    plt.tight_layout()
    return fig


def generate_excel(result_df, client_df, stats):
    """生成Excel文件"""
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        result_df.to_excel(writer, sheet_name='参会人员明细', index=False)

        summary = pd.DataFrame({
            '指标': ['总参会人数', 'A级客户', 'B级客户', 'AB级合计', '其他', 
                    'AB级覆盖率', 'A级覆盖率', 'B级覆盖率'],
            '数值': [
                stats['total_attendees'], stats['a_attendees'], stats['b_attendees'],
                stats['ab_attendees'], stats['other_attendees'],
                f"{stats['coverage_ab']}%", f"{stats['coverage_a']}%", f"{stats['coverage_b']}%",
            ]
        })
        summary.to_excel(writer, sheet_name='统计汇总', index=False)

        a_df = result_df[result_df['客户分级'] == 'A']
        a_df.to_excel(writer, sheet_name='A级客户', index=False)

        b_df = result_df[result_df['客户分级'] == 'B']
        b_df.to_excel(writer, sheet_name='B级客户', index=False)

        all_ab = set(client_df[client_df['客户分级'].isin(['A','B'])]['医生姓名'])
        attend = set(result_df[result_df['客户分级'].isin(['A','B'])]['姓名'])
        not_attend = all_ab - attend
        not_df = client_df[client_df['医生姓名'].isin(not_attend)]
        not_df.to_excel(writer, sheet_name='未参会AB客户', index=False)

    output.seek(0)
    return output.getvalue()

def generate_ppt(result_df, client_df, stats, meeting_name, meeting_date, fig1, fig2):
    """生成PPT文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    tmp_dir = tempfile.gettempdir()
    fig1_path = os.path.join(tmp_dir, 'chart1_temp.png')
    fig2_path = os.path.join(tmp_dir, 'chart2_temp.png')
    fig1.savefig(fig1_path, bbox_inches='tight', facecolor='white')
    fig2.savefig(fig2_path, bbox_inches='tight', facecolor='white')

    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
        bg.line.fill.background()

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.6))
            tf = tb2.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p.alignment = PP_ALIGN.CENTER
        return slide

    def add_content_slide(title, image_path=None, text_content=None, left_image=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
        bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

        if image_path and os.path.exists(image_path):
            if left_image:
                slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.9), width=Inches(7.5))
            else:
                slide.shapes.add_picture(image_path, Inches(0.3), Inches(0.9), width=Inches(12.7))

        if text_content:
            tb = slide.shapes.add_textbox(Inches(8.2) if left_image else Inches(0.5), 
                                          Inches(1.0) if left_image else Inches(6.5), 
                                          Inches(5.0) if left_image else Inches(12.3), 
                                          Inches(5.5) if left_image else Inches(0.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(text_content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14 if left_image else 13)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_before = Pt(6)
        return slide

    # 1. 封面
    add_title_slide(meeting_name, f"会议复盘报告 | {meeting_date}")

    # 2. 核心数据
    add_content_slide("核心数据总览", fig1_path, [
        f"总参会人数：{stats['total_attendees']} 人",
        f"A级客户：{stats['a_attendees']} 人（覆盖率 {stats['coverage_a']}%）",
        f"B级客户：{stats['b_attendees']} 人（覆盖率 {stats['coverage_b']}%）",
        f"AB级覆盖率：{stats['coverage_ab']}%（目标>=50%）",
        "",
        "诊断：覆盖率严重不足，会前邀约机制需优化"
    ], left_image=True)

    # 3. 深度分析
    add_content_slide("深度分析：覆盖率与竞品", fig2_path, [
        "各团队覆盖率差异巨大",
        "竞品高度集中（坦索罗辛为主）",
        "未参会客户数量庞大",
        "",
        "建议：标杆经验推广+低覆盖团队专项辅导"
    ], left_image=True)

    # 4. 问题诊断
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "问题诊断与改进建议"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x23, 0x7e)

    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.9), Inches(6.0), Inches(6.2))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
    left_bg.line.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5.6), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心问题"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    problems = [
        f"1. 覆盖率仅{stats['coverage_ab']}%，80%+核心客户未触达",
        "2. 部分团队覆盖率极低",
        "3. 40%+参会者观看<5分钟，打卡式参会严重",
        "4. 省级医院专家参与度低",
        "5. 主任医师占比不足"
    ]
    for prob in problems:
        p = tf.add_paragraph()
        p.text = prob
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)

    right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(0.9), Inches(6.3), Inches(6.2))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    right_bg.line.color.rgb = RGBColor(0x27, 0xAE, 0x60)
    tb = slide.shapes.add_textbox(Inches(6.9), Inches(1.0), Inches(5.9), Inches(6.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "改进方向"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
    solutions = [
        "1. 建立会前邀约KPI（A级>=60%，B级>=40%）",
        "2. 标杆经验推广（高覆盖团队最佳实践）",
        "3. 会中互动签到+实名准入机制",
        "4. 省级医院改用科室会形式",
        "5. 会后3天内未参会客户一对一跟进"
    ]
    for sol in solutions:
        p = tf.add_paragraph()
        p.text = sol
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_before = Pt(12)

    # 5. 总结
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x7e)
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.8))
    tf = tb2.text_frame
    p = tf.paragraphs[0]
    p.text = f"下次会议目标：AB级覆盖率 >= 40% | A级覆盖率 >= 50%"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    for p in [fig1_path, fig2_path]:
        if os.path.exists(p):
            os.remove(p)

    return output.getvalue()


# ========== 历史对比图表（Plotly，手机友好） ==========

def render_history_comparison():
    """渲染历史对比页面"""
    history = load_history()

    if len(history) == 0:
        st.info("📭 暂无历史记录。请先完成至少一次会议分析，数据会自动保存到这里。")
        return

    st.markdown(f"### 📊 历史对比分析（共 {len(history)} 次会议）")

    # 选择对比的会议
    meeting_options = [f"{h['meeting_date']} - {h['meeting_name']}" for h in history]
    selected = st.multiselect(
        "选择要对比的会议（建议选2-5次）",
        options=meeting_options,
        default=meeting_options[-3:] if len(meeting_options) >= 3 else meeting_options
    )

    if not selected:
        st.warning("请至少选择一次会议进行对比")
        return

    selected_indices = [meeting_options.index(s) for s in selected]
    selected_history = [history[i] for i in selected_indices]

    # 数据准备
    dates = [h['meeting_date'] for h in selected_history]
    names = [h['meeting_name'][:15] + '...' if len(h['meeting_name']) > 15 else h['meeting_name'] for h in selected_history]
    labels = [f"{d}<br>{n}" for d, n in zip(dates, names)]

    # 1. 覆盖率趋势
    fig1 = go.Figure()
    fig1.add_trace(go.Scatter(
        x=labels, y=[h['stats']['coverage_ab'] for h in selected_history],
        mode='lines+markers+text',
        name='AB级覆盖率',
        line=dict(color='#E74C3C', width=3),
        marker=dict(size=12),
        text=[f"{h['stats']['coverage_ab']}%" for h in selected_history],
        textposition='top center'
    ))
    fig1.add_trace(go.Scatter(
        x=labels, y=[h['stats']['coverage_a'] for h in selected_history],
        mode='lines+markers',
        name='A级覆盖率',
        line=dict(color='#FF6B6B', width=2, dash='dash'),
        marker=dict(size=8)
    ))
    fig1.add_trace(go.Scatter(
        x=labels, y=[h['stats']['coverage_b'] for h in selected_history],
        mode='lines+markers',
        name='B级覆盖率',
        line=dict(color='#4ECDC4', width=2, dash='dash'),
        marker=dict(size=8)
    ))
    fig1.add_hline(y=50, line_dash="dot", line_color="green", annotation_text="目标线50%")
    fig1.update_layout(
        title='📈 覆盖率变化趋势',
        xaxis_title='会议',
        yaxis_title='覆盖率 (%)',
        hovermode='x unified',
        height=400,
        margin=dict(l=40, r=40, t=60, b=80)
    )
    st.plotly_chart(fig1, use_container_width=True)

    # 2. 参会人数变化
    fig2 = make_subplots(rows=1, cols=2, subplot_titles=('总参会人数', 'AB级参会人数'))

    fig2.add_trace(go.Bar(
        x=labels, y=[h['stats']['total_attendees'] for h in selected_history],
        name='总参会', marker_color='#3498DB',
        text=[h['stats']['total_attendees'] for h in selected_history],
        textposition='auto'
    ), row=1, col=1)

    fig2.add_trace(go.Bar(
        x=labels, y=[h['stats']['ab_attendees'] for h in selected_history],
        name='AB级', marker_color='#27AE60',
        text=[h['stats']['ab_attendees'] for h in selected_history],
        textposition='auto'
    ), row=1, col=2)

    fig2.update_layout(height=350, showlegend=False, margin=dict(l=40, r=40, t=60, b=80))
    st.plotly_chart(fig2, use_container_width=True)

    # 3. 团队覆盖率对比（热力图）
    all_teams = set()
    for h in selected_history:
        all_teams.update(h.get('team_stats', {}).keys())
    all_teams = sorted(list(all_teams))

    if all_teams:
        coverage_matrix = []
        for h in selected_history:
            row = []
            for team in all_teams:
                row.append(h.get('team_stats', {}).get(team, {}).get('coverage', 0))
            coverage_matrix.append(row)

        fig3 = go.Figure(data=go.Heatmap(
            z=coverage_matrix,
            x=all_teams,
            y=[f"{h['meeting_date']}" for h in selected_history],
            colorscale=[[0, '#FDEDEC'], [0.3, '#FADBD8'], [0.5, '#F39C12'], [0.7, '#2ECC71'], [1, '#27AE60']],
            text=[[f"{v:.1f}%" for v in row] for row in coverage_matrix],
            texttemplate='%{text}',
            textfont={"size": 12},
            hovertemplate='团队: %{x}<br>会议: %{y}<br>覆盖率: %{z:.1f}%<extra></extra>'
        ))
        fig3.update_layout(
            title='🔥 各团队覆盖率对比热力图',
            height=300 + len(selected_history) * 40,
            xaxis_title='团队',
            yaxis_title='会议',
            margin=dict(l=40, r=40, t=60, b=60)
        )
        st.plotly_chart(fig3, use_container_width=True)

    # 4. 参会质量变化
    fig4 = go.Figure()
    fig4.add_trace(go.Bar(
        x=labels, y=[h['quality']['high'] for h in selected_history],
        name='高质量(>30分)', marker_color='#27AE60'
    ))
    fig4.add_trace(go.Bar(
        x=labels, y=[h['quality']['medium'] for h in selected_history],
        name='中质量(5-30分)', marker_color='#F39C12'
    ))
    fig4.add_trace(go.Bar(
        x=labels, y=[h['quality']['low'] for h in selected_history],
        name='低质量(<5分)', marker_color='#E74C3C'
    ))
    fig4.update_layout(
        title='📊 参会质量变化',
        barmode='group',
        height=400,
        xaxis_title='会议',
        yaxis_title='人数',
        margin=dict(l=40, r=40, t=60, b=80)
    )
    st.plotly_chart(fig4, use_container_width=True)

    # 历史记录表格
    st.markdown("### 📝 历史记录明细")
    history_df = pd.DataFrame([
        {
            '日期': h['meeting_date'],
            '会议名称': h['meeting_name'],
            '总参会': h['stats']['total_attendees'],
            'A级': h['stats']['a_attendees'],
            'B级': h['stats']['b_attendees'],
            'AB覆盖率': f"{h['stats']['coverage_ab']}%",
            '高质量': h['quality']['high'],
            '低质量': h['quality']['low']
        }
        for h in history
    ])
    st.dataframe(history_df, use_container_width=True, height=300)

    # 删除历史记录
    st.markdown("### ⚠️ 管理历史记录")
    col1, col2 = st.columns(2)
    with col1:
        delete_idx = st.selectbox("选择要删除的记录", range(len(history)), 
                                   format_func=lambda i: f"{history[i]['meeting_date']} - {history[i]['meeting_name']}")
    with col2:
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🗑️ 删除选中记录", use_container_width=True):
            if delete_history(delete_idx):
                st.success("已删除！请刷新页面查看更新")
                st.rerun()

    if st.button("🗑️ 清空所有历史记录", use_container_width=True):
        if os.path.exists(HISTORY_FILE):
            os.remove(HISTORY_FILE)
        st.success("已清空所有历史记录！")
        st.rerun()


# ========== 页面主体 ==========

st.markdown('<div class="main-title">📊 会议复盘分析自动化工具1.0</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">上传会议记录 + 客户总表 → 自动生成图表 + Excel + PPT</div>', unsafe_allow_html=True)

st.markdown("---")

# 文件上传区域（手机友好：单列布局）
watch_file = st.file_uploader("📁 上传会议观看记录（Excel）", type=['xlsx'], key='watch')
if watch_file:
    st.success(f"✅ 已上传：{watch_file.name}")

client_file = st.file_uploader("📁 上传AB级客户总表（Excel）", type=['xlsx'], key='client')
if client_file:
    st.success(f"✅ 已上传：{client_file.name}")

# 会议信息
col3, col4 = st.columns(2)
with col3:
    meeting_name = st.text_input("📝 会议名称", value="BPH/LUTS规范化诊疗区域会")
with col4:
    meeting_date = st.date_input("📅 会议日期", value=datetime.now())

st.markdown("---")

# 分析按钮
analyze_clicked = st.button("🚀 开始分析", use_container_width=True)

# 执行分析
if analyze_clicked:
    if not watch_file or not client_file:
        st.error("❌ 请先上传两个Excel文件！")
    else:
        with st.spinner("正在分析数据，请稍候..."):
            try:
                result = analyze_meeting(
                    watch_file.getvalue(),
                    client_file.getvalue(),
                    meeting_name,
                    meeting_date.strftime("%Y-%m-%d")
                )

                st.session_state['result'] = result
                st.success("✅ 分析完成！结果已保存到历史记录")
            except Exception as e:
                st.error(f"❌ 分析出错：{str(e)}")
                st.exception(e)

# 显示结果
if 'result' in st.session_state:
    result = st.session_state['result']
    stats = result['stats']
    result_df = result['result_df']

    st.markdown("---")
    st.markdown("### 📈 分析结果")

    # 关键指标卡片（手机：单列或双列）
    c1, c2, c3, c4, c5 = st.columns(5)

    with c1:
        st.markdown(f'<div class="metric-card" style="background: linear-gradient(135deg, #3498db 0%, #2980b9 100%);"><div class="metric-value">{stats["total_attendees"]}</div><div class="metric-label">总参会人数</div></div>', unsafe_allow_html=True)

    with c2:
        st.markdown(f'<div class="metric-card" style="background: linear-gradient(135deg, #e74c3c 0%, #c0392b 100%);"><div class="metric-value">{stats["a_attendees"]}</div><div class="metric-label">A级客户 ({stats["coverage_a"]}%)</div></div>', unsafe_allow_html=True)

    with c3:
        st.markdown(f'<div class="metric-card" style="background: linear-gradient(135deg, #f39c12 0%, #e67e22 100%);"><div class="metric-value">{stats["b_attendees"]}</div><div class="metric-label">B级客户 ({stats["coverage_b"]}%)</div></div>', unsafe_allow_html=True)

    with c4:
        st.markdown(f'<div class="metric-card" style="background: linear-gradient(135deg, #27ae60 0%, #2ecc71 100%);"><div class="metric-value">{stats["ab_attendees"]}</div><div class="metric-label">AB级合计</div></div>', unsafe_allow_html=True)

    with c5:
        cov_color = "#27ae60" if stats['coverage_ab'] >= 50 else "#e74c3c"
        st.markdown(f'<div class="metric-card" style="background: linear-gradient(135deg, {cov_color} 0%, #c0392b 100%);"><div class="metric-value">{stats["coverage_ab"]}%</div><div class="metric-label">AB级覆盖率</div></div>', unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # 标签页（手机友好）
    tab1, tab2, tab3, tab4, tab5 = st.tabs(["📊 数据分析", "🔍 深度分析", "📋 参会明细", "📈 历史对比", "⬇️ 下载报告"])

    with tab1:
        st.pyplot(result['fig1'], use_container_width=True)

    with tab2:
        st.pyplot(result['fig2'], use_container_width=True)

    with tab3:
        filter_grade = st.multiselect("筛选客户分级", ['A', 'B', '未匹配'], default=['A', 'B', '未匹配'])
        filtered_df = result_df[result_df['客户分级'].isin(filter_grade)]
        st.dataframe(filtered_df, use_container_width=True, height=500)
        st.caption(f"显示 {len(filtered_df)} / {len(result_df)} 条记录")

    with tab4:
        render_history_comparison()

    with tab5:
        col_d1, col_d2 = st.columns(2)

        with col_d1:
            st.download_button(
                label="📥 下载Excel明细",
                data=result['excel_bytes'],
                file_name=f"{meeting_date.strftime('%Y%m%d')}_{meeting_name}_分析明细.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )

        with col_d2:
            st.download_button(
                label="📥 下载PPT报告",
                data=result['ppt_bytes'],
                file_name=f"{meeting_date.strftime('%Y%m%d')}_{meeting_name}_复盘报告.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

        st.markdown('<div class="info-box">Excel包含5个工作表：参会明细、统计汇总、A级客户、B级客户、未参会AB客户</div>', unsafe_allow_html=True)
        st.markdown('<div class="info-box">PPT包含5页：封面、核心数据、深度分析、问题诊断、总结</div>', unsafe_allow_html=True)

# 页脚
st.markdown("---")
st.markdown("<div style='text-align:center;color:#999;font-size:12px;'>会议复盘分析自动化工具 | 安徽大区-王斌 | v2.0</div>", unsafe_allow_html=True)
